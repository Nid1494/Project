import pandas as pd
import numpy as np
from scipy.stats import pearsonr
from sklearn.linear_model import LinearRegression
from sklearn.ensemble import RandomForestRegressor
from sklearn.neural_network import MLPRegressor
from sklearn.cross_decomposition import PLSRegression
from sklearn.metrics import mean_squared_error
from sklearn.model_selection import train_test_split
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches


def generate_time_series_plot(actual_values, predicted_values, model_name, filename):
    fig, ax = plt.subplots(figsize=(10, 6))

    time_values = np.arange(1, len(actual_values) + 1)

    # Convert actual_values and predicted_values to NumPy arrays
    actual_values_np = np.array(actual_values)
    predicted_values_np = np.array(predicted_values)

    # Identify outliers (replace this with your own outlier detection method)
    outlier_indices = np.where(np.abs(actual_values_np - predicted_values_np) > 50.0)[0]

    ax.plot(time_values, actual_values_np, label='Actual', color='blue')
    ax.plot(time_values, predicted_values_np, label='Predicted', color='orange')

    # Mark outliers with red crosses
    ax.scatter(time_values[outlier_indices], predicted_values_np[outlier_indices], marker='x', color='red', label='Outliers')

    ax.set_xlabel('Time')
    ax.set_ylabel('Value')
    ax.set_title(f"{model_name} - Actual vs Predicted")
    ax.legend()
    plt.tight_layout()

    fig.savefig(filename)
    plt.close(fig)

    return fig


def generate_bar_graph(variable_importance, variable_names, model_name):
    plt.figure(figsize=(10, 6))

    # Sort variable importance in descending order
    sorted_indices = variable_importance.argsort()[::-1]
    sorted_importance = variable_importance[sorted_indices]
    sorted_names = variable_names[sorted_indices]

    plt.barh(range(len(variable_names)), sorted_importance, tick_label=sorted_names)
    plt.xlabel('Importance/Weight')
    plt.ylabel('Variable')
    plt.title(f"{model_name} - Variable Importance/Parameter Weights")
    plt.axvline(x=0, color='black', linewidth=0.5)
    plt.tight_layout()

    # Save the bar graph
    plt.savefig(f"{model_name}_variable_importance.png")
    plt.close()


# Load the dataset from XLSX file
data = pd.read_excel('PNS_DataForModeling_Project.xlsx', skiprows=2)
data = data.dropna(axis=1, how='all')  # Remove columns with all NaN values

# Remove the rows used as column names
data = data.iloc[2:].reset_index(drop=True)

# Replace '--' values with NaN
data = data.replace('--', np.nan)

# Convert the first column (Sampling Time) to datetime type with correct format
data['Sampling Time'] = pd.to_datetime(data['Sampling Time'], format='%d/%m/%y %H:%M')

# Set the 'Sampling Time' column as the index
data.set_index('Sampling Time', inplace=True)

# Replace 0 values with NaN before forward filling
data.replace(0, np.nan, inplace=True)

# Select columns for which you want to perform interpolation and resampling
columns_to_interpolate = data.columns[::-1]  # Exclude 'Sampling Time' and last column

# Interpolate missing values for selected columns
data[columns_to_interpolate] = data[columns_to_interpolate].interpolate(method='linear')


# Resample the data to hourly intervals and calculate the meanwhile ignoring missing values
def mean_without_nan(x):
    if np.isnan(x).all():
        return np.nan
    return np.nanmean(x)


data = data.resample('H').apply(mean_without_nan)

sampling_time = data.index

# Compute correlation
target_column_name = "Finish Mill 4 Blaine"
blaine_corr = data.corr()[target_column_name].abs().sort_values(ascending=False)
top_variables = blaine_corr.drop(target_column_name).nlargest(15)
blaine_corr_filtered = top_variables.dropna()  # Filter out columns with NaN correlation values

# Model training and evaluation
print("Root Mean Squared Error (RMSE)")

# Multi-Variate Linear Regression (MLR)
X = data[top_variables.index]  # Independent variables
y = data[target_column_name]  # Dependent variable
X_train_mlr, X_test_mlr, y_train_mlr, y_test_mlr = train_test_split(X, y, test_size=0.2, random_state=42)
mlr_model = LinearRegression()
mlr_model.fit(X_train_mlr, y_train_mlr)
mlr_y_pred = mlr_model.predict(X_test_mlr)
mlr_train_pred = mlr_model.predict(X_train_mlr)
mlr_corr_coef, _ = pearsonr(y_test_mlr, mlr_y_pred)  # Calculate correlation coefficients
mlr_mse = mean_squared_error(y_test_mlr, mlr_y_pred)
mlr_rmse = np.sqrt(mean_squared_error(y_test_mlr, mlr_y_pred))
print("\nMulti-Variate Linear Regression (MLR) - Root Mean Squared Error (RMSE):", mlr_rmse)

# Random Forest
X = data[top_variables.index]  # Independent variables
y = data[target_column_name]  # Dependent variable
X_train_rf, X_test_rf, y_train_rf, y_test_rf = train_test_split(X, y, test_size=0.2, random_state=42)
rf_model = RandomForestRegressor(n_estimators=100, random_state=42)
rf_model.fit(X_train_rf, y_train_rf)
rf_y_pred = rf_model.predict(X_test_rf)
rf_train_pred = mlr_model.predict(X_train_rf)
rf_corr_coef, _ = pearsonr(y_test_rf, rf_y_pred)  # Calculate correlation coefficients
rf_mse = mean_squared_error(y_test_rf, rf_y_pred)
rf_rmse = np.sqrt(mean_squared_error(y_test_rf, rf_y_pred))
print("\nRandom Forest - Root Mean Squared Error (RMSE):", rf_rmse)

# Partial Least Squares (PLS)
X = data[top_variables.index]  # Independent variables
y = data[target_column_name]  # Dependent variable
X_train_pls, X_test_pls, y_train_pls, y_test_pls = train_test_split(X, y, test_size=0.2, random_state=42)
n_components = min(15, X_train_pls.shape[1])  # Determine the number of components based on the input data
pls_model = PLSRegression(n_components=n_components)
pls_model.fit(X_train_pls, y_train_pls)
pls_y_pred = pls_model.predict(X_test_pls)
y_test_array = y_test_pls.to_numpy()
nan_indices = np.isnan(y_test_array.flatten()) | np.isnan(pls_y_pred.flatten())
y_test_clean = y_test_array[~nan_indices].astype(np.float64)
pls_y_pred_clean = pls_y_pred[~nan_indices].astype(np.float64)
y_test_clean = y_test_clean.reshape(-1, 1)  # Reshape the arrays to have the same shape
pls_y_pred_clean = pls_y_pred_clean.reshape(-1, 1)  # Reshape the arrays to have the same shape
corr_matrix = np.corrcoef(y_test_clean.flatten(), pls_y_pred_clean.flatten())  # Calculate the correlation matrix
pls_corr_coef = corr_matrix[0, 1]  # Extract the correlation coefficient
pls_mse = mean_squared_error(y_test_clean, pls_y_pred)
pls_rmse = np.sqrt(mean_squared_error(y_test_clean, pls_y_pred_clean))
pls_train_pred = mlr_model.predict(X_train_pls)
print("\nPartial Least Squares (PLS) - Root Mean Squared Error (RMSE):", pls_rmse)

# Multi-Layer Perceptron (MLP)
X = data[top_variables.index]  # Independent variables
y = data[target_column_name]  # Dependent variable
X_train_mlp, X_test_mlp, y_train_mlp, y_test_mlp = train_test_split(X, y, test_size=0.2, random_state=42)
mlp_model = MLPRegressor(hidden_layer_sizes=(100,), max_iter=10000)
mlp_model.fit(X_train_mlp, y_train_mlp)
mlp_y_pred = mlp_model.predict(X_test_mlp)
mlp_train_pred = mlr_model.predict(X_train_mlp)
mlp_corr_coef, _ = pearsonr(y_test_mlp, mlp_y_pred)  # Calculate correlation coefficients
mlp_mse = mean_squared_error(y_test_mlp, mlp_y_pred)
mlp_rmse = np.sqrt(mean_squared_error(y_test_mlp, mlp_y_pred))
print("\nMulti-Layer Perceptron (MLP) - Root Mean Squared Error (RMSE):", mlp_rmse)

# Select top variables with correlation > 0.4 or < -0.4
print("\nVariable Influence and Correlation:")
top_variables = blaine_corr_filtered[(blaine_corr_filtered > 0.4) | (blaine_corr_filtered < -0.4)].nlargest(15)
for variable, correlation in top_variables.items():
    print(f"{variable: <30} = {correlation: .8f}")

# Data Summary to calculate min, max, mean and std
print("\nData Summary:")
data_summary = data.iloc[:, 1:].describe().loc[['min', 'max', 'mean', 'std']]
print(data_summary)

# Create a new PowerPoint presentation
prs = Presentation()

# Slide layout for content and two content objects (placeholders)
slide_layout = prs.slide_layouts[5]

# Slide 1: Performance Matrix
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

# Define the table data for the Performance Matrix slide
metrics = [
    ['Model', 'RMSE', 'Correlation Coefficient', 'MSE'],
    ['Multi-Variate Linear Regression (MLR)', mlr_rmse, mlr_corr_coef, mlr_mse],
    ['Random Forest', rf_rmse, rf_corr_coef, rf_mse],
    ['Partial Least Squares (PLS)', pls_rmse, pls_corr_coef, pls_mse],
    ['Multi-Layer Perceptron (MLP)', mlp_rmse, mlp_corr_coef, mlp_mse]
]

left = Inches(1)
top = Inches(1.5)
width = Inches(6)
height = Inches(2.5)

# Add the table to the slide
table = shapes.add_table(rows=len(metrics), cols=len(metrics[0]), left=left, top=top, width=width, height=height).table

# Populate the table with data
for i, row in enumerate(metrics):
    for j, value in enumerate(row):
        table.cell(i, j).text = str(value)
# Slide 2: MLR Time Series Plot
slide = prs.slides.add_slide(slide_layout)

left = Inches(1)
top = Inches(1.5)
width = Inches(6)
height = Inches(4)

# Add MLR time series plot
mlr_time_series_fig = generate_time_series_plot(y_test_mlr.to_numpy(), mlr_y_pred,
                                                'Multi-Variate Linear Regression (MLR) - Testing Set',
                                                'mlr_testing_plot.png')
mlr_time_series_fig.savefig('mlr_testing_plot.png')
# Generate time series plot for Actual Output Variable trend in Training Set - MLR
generate_time_series_plot(y_train_mlr, mlr_model.predict(X_train_mlr),
                          'Multi-Variate Linear Regression (MLR) - Training Set', 'mlr_training_plot.png')
mlr_training_time_series_fig = generate_time_series_plot(y_train_mlr.to_numpy(), mlr_train_pred,
                                                         'Multi-Variate Linear Regression (MLR) - Training Set',
                                                         'mlr_training_plot.png')
mlr_training_time_series_fig.savefig('mlr_training_plot.png')
slide.shapes.add_picture('mlr_training_plot.png', left, top, width, height)
slide.shapes.add_picture('mlr_testing_plot.png', left, top, width, height)

# Slide 3: RF Time Series Plot
slide = prs.slides.add_slide(slide_layout)

# Add RF time series plot
rf_time_series_fig = generate_time_series_plot(y_test_rf.to_numpy(), rf_y_pred, 'Random Forest - Testing Set',
                                               'rf_testing_plot.png')
rf_time_series_fig.savefig('rf_testing_plot.png')
generate_time_series_plot(y_train_rf, rf_model.predict(X_train_rf), 'Random Forest - Training Set',
                          'rf_training_plot.png')
rf_training_time_series_fig = generate_time_series_plot(y_train_rf.to_numpy(), rf_train_pred,
                                                        'Random Forest - Training Set', 'rf_training_plot.png')
rf_training_time_series_fig.savefig('rf_training_plot.png')
slide.shapes.add_picture('rf_training_plot.png', left, top, width, height)
slide.shapes.add_picture('rf_testing_plot.png', left + 2 * width, top, width, height)

# Slide 4: PLS Time Series Plot
slide = prs.slides.add_slide(slide_layout)

# Add PLS time series plot
pls_time_series_fig = generate_time_series_plot(y_test_clean, pls_y_pred_clean,
                                                'Partial Least Squares (PLS) - Testing Set', 'pls_testing_plot.png')
pls_time_series_fig.savefig('pls_testing_plot.png')
generate_time_series_plot(y_train_pls, pls_model.predict(X_train_pls),
                          'Partial Least Squares (PLS) - Training Set', 'pls_training_plot.png')
pls_training_time_series_fig = generate_time_series_plot(y_train_pls.to_numpy(), pls_train_pred,
                                                         'Partial Least Squares (PLS) - Training Set',
                                                         'pls_training_plot.png')
pls_training_time_series_fig.savefig('pls_training_plot.png')
slide.shapes.add_picture('pls_training_plot.png', left, top, width, height)
slide.shapes.add_picture('pls_testing_plot.png', left, top + height, width, height)

# Slide 5: MLP Time Series Plot
slide = prs.slides.add_slide(slide_layout)

# Add MLP time series plot
mlp_time_series_fig = generate_time_series_plot(y_test_mlp.to_numpy(), mlp_y_pred,
                                                'Multi-Layer Perceptron (MLP) - Testing Set', 'mlp_testing_plot.png')
mlp_time_series_fig.savefig('mlp_testing_plot.png')
generate_time_series_plot(y_train_mlp, mlp_model.predict(X_train_mlp),
                          'Multi-Layer Perceptron (MLP) - Training Set', 'mlp_training_plot.png')
mlp_training_time_series_fig = generate_time_series_plot(y_train_mlp.to_numpy(), mlp_train_pred,
                                                         'Multi-Layer Perceptron (MLP) - Training Set',
                                                         'mlp_training_plot.png')
mlp_training_time_series_fig.savefig('mlp_training_plot.png')
slide.shapes.add_picture('mlp_training_plot.png', left, top + height, width, height)
slide.shapes.add_picture('mlp_testing_plot.png', left + 2 * width, top, width, height)

# Multi-Variate Linear Regression (MLR)
mlr_variable_importance = mlr_model.coef_  # Get the coefficient values
mlr_variable_names = X.columns  # Get the names of the independent variables
generate_bar_graph(mlr_variable_importance, mlr_variable_names, "MLR Variable Importance")

# Random Forest
rf_variable_importance = rf_model.feature_importances_  # Get the feature importance values
rf_variable_names = X.columns  # Get the names of the independent variables
generate_bar_graph(rf_variable_importance, rf_variable_names, "Random Forest Variable Importance")

# Partial Least Squares (PLS)
pls_variable_importance = np.abs(pls_model.coef_)  # Get the coefficient values and take absolute values
pls_variable_importance_avg = np.mean(pls_variable_importance, axis=0)  # Average importance across components
pls_variable_names = X.columns  # Get the names of the independent variables
generate_bar_graph(pls_variable_importance_avg, pls_variable_names, "PLS Variable Importance")

# Define the model names
model_names = [
    'Multi-Variate Linear Regression (MLR)',
    'Random Forest',
    'Partial Least Squares (PLS)',
    'Multi-Layer Perceptron (MLP)'
]

# Iterate over the slides and add the model names
for idx, slide in enumerate(prs.slides):
    if 0 < idx <= len(model_names):
        model = model_names[idx - 1]
        slide.placeholders[0].text = model  # Replace 0 with the correct index of the placeholder

#  Add Variable Influence and its correlation
slide_layout = prs.slide_layouts[5]  # Use layout for a title and content slide
slide = prs.slides.add_slide(slide_layout)
# Add slide with table
table = slide.shapes.add_table(rows=len(top_variables) + 1,
                               cols=3, left=left, top=top, width=width, height=height).table

# Add title
titles = slide.shapes.title
titles.text = "Variables Shortlisted based on Correlation"

# Create table shape
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(4)

# Set column names
table.cell(0, 0).text = "Variable"
table.cell(0, 1).text = "Correlation"
table.cell(0, 2).text = "Correlation Value"

# Populate table with data
for i, variable in enumerate(top_variables.index):
    table.cell(i + 1, 0).text = variable
    table.cell(i + 1, 1).text = str(round(blaine_corr[variable], 2))
    table.cell(i + 1, 2).text = str(round(data.corr()[target_column_name][variable], 2))

# Save the presentation
prs.save('Performance_Matrix.pptx')

# Save the modified DataFrame to a new Excel file
with pd.ExcelWriter('Modified_PNS_Data.xlsx', engine='openpyxl') as writer:
    # Write the DataFrame to the Excel file
    data.to_excel(writer, sheet_name='Sheet1', index=True)  # Keep index column

    # Access the workbook and the sheet
    workbook = writer.book
    sheet = workbook['Sheet1']

    # Create a new DataFrame with only the sampling_time column
    sampling_time_df = pd.DataFrame({'Sampling Time': sampling_time})

    sheet.cell(row=1, column=1, value='Sampling Time')

    # Write the Sampling Time values to the first column
    for i, time in enumerate(sampling_time_df['Sampling Time'], start=2):
        sheet.cell(row=i, column=1, value=time)

    # Save the modified Excel file
    workbook.save('Modified_PNS_Data.xlsx')


# Create a Pandas DataFrame to store the results
results_df = pd.DataFrame({
    'Model': ['MLR', 'RF', 'PLS', 'MLP'],
    'Actual': [y_test_mlr, y_test_rf, y_test_clean, y_test_mlp],
    'Predicted': [mlr_y_pred, rf_y_pred, pls_y_pred_clean, mlp_y_pred],
    'Train Actual': [y_train_mlr, y_train_rf, y_train_pls, y_train_mlp],
    'Train Predicted': [mlr_model.predict(X_train_mlr), rf_model.predict(X_train_rf),
                        pls_model.predict(X_train_pls), mlp_model.predict(X_train_mlp)]
})

# Convert the string representation back to numeric values
results_df['Actual'] = results_df['Actual'].apply(lambda x: [float(val) for val in x])
results_df['Predicted'] = results_df['Predicted'].apply(lambda x: [float(val) for val in x])
results_df['Train Actual'] = results_df['Train Actual'].apply(lambda x: [float(val) for val in x])
results_df['Train Predicted'] = results_df['Train Predicted'].apply(lambda x: [float(val) for val in x])

# For MLR
mlr_feature_names = X.columns.tolist()  # Assuming X is your feature data
mlr_coefs = mlr_model.coef_
mlr_coefs_abs_sum = np.abs(mlr_coefs).sum(axis=0)  # Sum along axis 0

# For PLS
pls_feature_names = X_train_pls.columns.tolist()  # Assuming X_train_pls is your feature data
pls_loadings = pls_model.x_loadings_
pls_loadings_abs_sum = np.abs(pls_loadings).sum(axis=0)  # Sum along axis 0

# Visualize Variable Importance (Loading Factors) - Separate Plots
plt.figure(figsize=(12, 8))

# MLR
plt.subplot(2, 1, 1)
plt.bar(X_train_mlr.columns, mlr_coefs_abs_sum)
plt.xlabel('Features')
plt.ylabel('Variable Importance (Parameter Weight) - MLR')
plt.xticks(rotation=45)
plt.tight_layout()

# PLS
plt.subplot(2, 1, 2)
plt.bar(X_train_pls.columns, pls_loadings_abs_sum)
plt.xlabel('Features')
plt.ylabel('Variable Importance (Loading Factor) - PLS')
plt.xticks(rotation=45)
plt.tight_layout()

plt.show()

# Save the DataFrame to a new Excel file with separate sheets for each model
with pd.ExcelWriter('Model_Predictions.xlsx', engine='xlsxwriter') as excel_writer:
    for models_name in results_df['Model']:
        df = results_df[results_df['Model'] == models_name]

        # Convert the DataFrame to a list of lists
        df_values = df.values.tolist()

        # Create a new sheet for the model
        sheet_name = models_name
        df.to_excel(excel_writer, sheet_name=sheet_name, index=False)

        # Get the xlsxwriter workbook and worksheet objects
        workbook = excel_writer.book
        worksheet = excel_writer.sheets[sheet_name]

        # Create a chart object
        chart = workbook.add_chart({'type': 'scatter'})

        # Configure the first series
        chart.add_series({
            'name': [sheet_name, 1, 3],
            'categories': [sheet_name, 2, 0, len(df_values), 0],
            'values': [sheet_name, 2, 3, len(df_values), 3],
        })

        # Add a chart title
        chart.set_title({'name': 'Actual vs Predicted'})

        # Insert the chart into the worksheet
        worksheet.insert_chart('G2', chart)
